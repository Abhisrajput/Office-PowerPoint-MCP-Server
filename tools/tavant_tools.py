"""
Tavant Status Report Tool - Custom extension for Office-PowerPoint-MCP-Server
Add this file to the tools/ directory of the cloned repository
"""

import json
import os
from datetime import datetime
from typing import Any, Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Tavant Colors
TAVANT_COLORS = {
    "primary": RgbColor(0, 0, 0),        # Black
    "accent": RgbColor(242, 101, 34),    # Orange #F26522
    "white": RgbColor(255, 255, 255),
    "dark": RgbColor(51, 51, 51),
    "light": RgbColor(245, 245, 245),
}


def create_tavant_status_report(
    project_name: str,
    period_label: str,
    accomplishments: List[str],
    priorities: List[Dict[str, str]],
    risks: List[Dict[str, str]],
    milestones: List[Dict[str, str]],
    upcoming_milestones: List[Dict[str, str]] = None,
    contact_info: str = "",
    output_path: str = None
) -> Dict[str, Any]:
    """
    Create a Tavant-branded Weekly Status Report presentation.
    
    Args:
        project_name: Name of the project
        period_label: Report period (e.g., "Week Ending Dec 13, 2025")
        accomplishments: List of accomplishment strings
        priorities: List of dicts with 'description' and 'owner' keys
        risks: List of dicts with 'description', 'owner', 'target_date', 'status' keys
        milestones: List of dicts with 'description', 'target_date', 'status' keys
        upcoming_milestones: List of dicts with 'description', 'target_date', 'owner' keys
        contact_info: Contact information for the thank you slide
        output_path: Path to save the presentation (optional)
    
    Returns:
        Dictionary with success status and file path
    """
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title Slide
        _create_title_slide(prs, project_name, period_label)
        
        # Slide 2: Executive Summary
        _create_executive_summary_slide(prs, accomplishments, priorities, risks)
        
        # Slide 3: Key Milestones
        _create_milestones_slide(prs, milestones, upcoming_milestones or [])
        
        # Slide 4: Thank You
        _create_thank_you_slide(prs, contact_info)
        
        # Save presentation
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"Tavant_WSR_{project_name.replace(' ', '_')}_{timestamp}.pptx"
        
        prs.save(output_path)
        
        return {
            "success": True,
            "file_path": output_path,
            "slides_created": 4,
            "message": f"Tavant Status Report created successfully"
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }


def _create_title_slide(prs: Presentation, project_name: str, period_label: str):
    """Create the title slide with Tavant branding."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Black background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = TAVANT_COLORS["primary"]
    background.line.fill.background()
    
    # Orange accent bar at bottom
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.33), Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TAVANT_COLORS["accent"]
    accent_bar.line.fill.background()
    
    # TAVANT brand name
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(4), Inches(0.5))
    brand_tf = brand_box.text_frame
    brand_p = brand_tf.paragraphs[0]
    brand_p.text = "TAVANT"
    brand_p.font.size = Pt(24)
    brand_p.font.bold = True
    brand_p.font.color.rgb = TAVANT_COLORS["accent"]
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "WEEKLY STATUS REPORT"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = TAVANT_COLORS["white"]
    
    # Project name
    project_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(10), Inches(0.5))
    project_tf = project_box.text_frame
    project_p = project_tf.paragraphs[0]
    project_p.text = f"Project: {project_name}"
    project_p.font.size = Pt(24)
    project_p.font.color.rgb = TAVANT_COLORS["white"]
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4), Inches(0.4))
    date_tf = date_box.text_frame
    date_p = date_tf.paragraphs[0]
    date_p.text = datetime.now().strftime("%B %d, %Y")
    date_p.font.size = Pt(18)
    date_p.font.color.rgb = TAVANT_COLORS["accent"]
    
    # Period label
    period_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(8), Inches(0.3))
    period_tf = period_box.text_frame
    period_p = period_tf.paragraphs[0]
    period_p.text = period_label
    period_p.font.size = Pt(14)
    period_p.font.color.rgb = TAVANT_COLORS["white"]


def _create_executive_summary_slide(
    prs: Presentation, 
    accomplishments: List[str], 
    priorities: List[Dict[str, str]], 
    risks: List[Dict[str, str]]
):
    """Create the executive summary slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.36), Inches(0.3), Inches(12), Inches(0.5))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Executive Summary"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = TAVANT_COLORS["dark"]
    
    # Orange underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(0.8), Inches(12.6), Inches(0.05)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = TAVANT_COLORS["accent"]
    underline.line.fill.background()
    
    # Key Accomplishments header
    _add_section_header(slide, "Key Accomplishments for Last Period", 0.36, 1.0, 5.8)
    
    # Accomplishments content
    acc_box = slide.shapes.add_textbox(Inches(0.36), Inches(1.4), Inches(5.8), Inches(2.0))
    acc_tf = acc_box.text_frame
    acc_tf.word_wrap = True
    for i, acc in enumerate(accomplishments[:5]):
        p = acc_tf.paragraphs[0] if i == 0 else acc_tf.add_paragraph()
        p.text = f"• {acc}"
        p.font.size = Pt(10)
        p.font.color.rgb = TAVANT_COLORS["dark"]
        p.space_after = Pt(4)
    
    # Top Priorities header
    _add_section_header(slide, "Top Priorities for Next Period", 6.5, 1.0, 6.5)
    
    # Priorities table
    _add_priorities_table(slide, priorities)
    
    # Risks header
    _add_section_header(slide, "Key Risks, Issues and Action Items", 0.36, 3.6, 12.6)
    
    # Risks table
    _add_risks_table(slide, risks)
    
    # Footer
    _add_footer(slide)


def _create_milestones_slide(
    prs: Presentation, 
    milestones: List[Dict[str, str]], 
    upcoming: List[Dict[str, str]]
):
    """Create the milestones slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.36), Inches(0.3), Inches(12), Inches(0.5))
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_p.text = "Key Milestones"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = TAVANT_COLORS["dark"]
    
    # Orange underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(0.8), Inches(12.6), Inches(0.05)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = TAVANT_COLORS["accent"]
    underline.line.fill.background()
    
    # Key Milestones header
    _add_section_header(slide, "Key Milestones and Status", 0.36, 1.0, 12.6)
    
    # Milestones table
    _add_milestones_table(slide, milestones, 1.4)
    
    # Upcoming Milestones header
    _add_section_header(slide, "Upcoming Key Milestones", 0.36, 4.2, 12.6)
    
    # Upcoming table
    _add_upcoming_table(slide, upcoming, 4.6)
    
    # Footer
    _add_footer(slide)


def _create_thank_you_slide(prs: Presentation, contact_info: str):
    """Create the thank you slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Black background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = TAVANT_COLORS["primary"]
    background.line.fill.background()
    
    # Orange accent bar at bottom
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.33), Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = TAVANT_COLORS["accent"]
    accent_bar.line.fill.background()
    
    # THANK YOU text
    thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.33), Inches(1.2))
    thank_tf = thank_box.text_frame
    thank_p = thank_tf.paragraphs[0]
    thank_p.text = "THANK YOU"
    thank_p.font.size = Pt(56)
    thank_p.font.bold = True
    thank_p.font.color.rgb = TAVANT_COLORS["white"]
    thank_p.alignment = PP_ALIGN.CENTER
    
    # TAVANT brand
    brand_box = slide.shapes.add_textbox(Inches(8.5), Inches(5.5), Inches(4), Inches(0.5))
    brand_tf = brand_box.text_frame
    brand_p = brand_tf.paragraphs[0]
    brand_p.text = "TAVANT"
    brand_p.font.size = Pt(24)
    brand_p.font.bold = True
    brand_p.font.color.rgb = TAVANT_COLORS["accent"]
    brand_p.alignment = PP_ALIGN.RIGHT
    
    # Contact info
    if contact_info:
        contact_box = slide.shapes.add_textbox(Inches(8.5), Inches(6.1), Inches(4), Inches(0.3))
        contact_tf = contact_box.text_frame
        contact_p = contact_tf.paragraphs[0]
        contact_p.text = contact_info
        contact_p.font.size = Pt(10)
        contact_p.font.color.rgb = TAVANT_COLORS["white"]
        contact_p.alignment = PP_ALIGN.RIGHT


def _add_section_header(slide, text: str, left: float, top: float, width: float):
    """Add a section header with orange background."""
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.35)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = TAVANT_COLORS["accent"]
    header.line.fill.background()
    
    # Add text to the shape
    header.text_frame.paragraphs[0].text = text
    header.text_frame.paragraphs[0].font.size = Pt(11)
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["white"]


def _add_priorities_table(slide, priorities: List[Dict[str, str]]):
    """Add the priorities table."""
    rows = min(len(priorities), 5) + 1  # Header + data rows
    cols = 3
    
    table = slide.shapes.add_table(rows, cols, Inches(6.5), Inches(1.4), Inches(6.5), Inches(0.3 * rows)).table
    
    # Set column widths
    table.columns[0].width = Inches(0.4)
    table.columns[1].width = Inches(4.2)
    table.columns[2].width = Inches(1.9)
    
    # Header row
    headers = ["#", "Description", "Owner"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TAVANT_COLORS["accent"]
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["white"]
    
    # Data rows
    for i, priority in enumerate(priorities[:5]):
        table.cell(i + 1, 0).text = str(i + 1)
        table.cell(i + 1, 1).text = priority.get("description", "")
        table.cell(i + 1, 2).text = priority.get("owner", "")
        
        for j in range(3):
            table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["dark"]


def _add_risks_table(slide, risks: List[Dict[str, str]]):
    """Add the risks table."""
    rows = min(len(risks), 3) + 1  # Header + data rows
    cols = 5
    
    table = slide.shapes.add_table(rows, cols, Inches(0.36), Inches(4.0), Inches(12.6), Inches(0.3 * rows)).table
    
    # Set column widths
    table.columns[0].width = Inches(0.4)
    table.columns[1].width = Inches(6.5)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(2.2)
    
    # Header row
    headers = ["#", "Action Item", "Owner", "Target Date", "Status"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TAVANT_COLORS["accent"]
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["white"]
    
    # Data rows
    for i, risk in enumerate(risks[:3]):
        table.cell(i + 1, 0).text = str(i + 1)
        table.cell(i + 1, 1).text = risk.get("description", "")
        table.cell(i + 1, 2).text = risk.get("owner", "")
        table.cell(i + 1, 3).text = risk.get("target_date", "")
        table.cell(i + 1, 4).text = risk.get("status", "")
        
        for j in range(5):
            table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["dark"]


def _add_milestones_table(slide, milestones: List[Dict[str, str]], top: float):
    """Add the milestones table."""
    rows = min(len(milestones), 6) + 1
    cols = 4
    
    table = slide.shapes.add_table(rows, cols, Inches(0.36), Inches(top), Inches(12.6), Inches(0.3 * rows)).table
    
    # Set column widths
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(7.5)
    table.columns[2].width = Inches(2.3)
    table.columns[3].width = Inches(2.3)
    
    # Header row
    headers = ["#", "Milestone Description", "Target Date", "Status"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TAVANT_COLORS["accent"]
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["white"]
    
    # Data rows
    for i, milestone in enumerate(milestones[:6]):
        table.cell(i + 1, 0).text = str(i + 1)
        table.cell(i + 1, 1).text = milestone.get("description", "")
        table.cell(i + 1, 2).text = milestone.get("target_date", "")
        table.cell(i + 1, 3).text = milestone.get("status", "")
        
        for j in range(4):
            table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["dark"]


def _add_upcoming_table(slide, upcoming: List[Dict[str, str]], top: float):
    """Add the upcoming milestones table."""
    rows = min(len(upcoming), 3) + 1
    cols = 3
    
    table = slide.shapes.add_table(rows, cols, Inches(0.36), Inches(top), Inches(12.6), Inches(0.3 * rows)).table
    
    # Set column widths
    table.columns[0].width = Inches(8)
    table.columns[1].width = Inches(2.3)
    table.columns[2].width = Inches(2.3)
    
    # Header row
    headers = ["Milestone", "Target Date", "Owner"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TAVANT_COLORS["accent"]
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["white"]
    
    # Data rows
    for i, item in enumerate(upcoming[:3]):
        table.cell(i + 1, 0).text = item.get("description", "")
        table.cell(i + 1, 1).text = item.get("target_date", "")
        table.cell(i + 1, 2).text = item.get("owner", "")
        
        for j in range(3):
            table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(9)
            table.cell(i + 1, j).text_frame.paragraphs[0].font.color.rgb = TAVANT_COLORS["dark"]


def _add_footer(slide):
    """Add TAVANT footer to slide."""
    footer_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    footer_tf = footer_box.text_frame
    footer_p = footer_tf.paragraphs[0]
    footer_p.text = "TAVANT"
    footer_p.font.size = Pt(10)
    footer_p.font.bold = True
    footer_p.font.color.rgb = TAVANT_COLORS["accent"]
    footer_p.alignment = PP_ALIGN.RIGHT


# MCP Tool Registration (to be added to ppt_mcp_server.py)
TAVANT_TOOL_DEFINITION = {
    "name": "create_tavant_status_report",
    "description": "Create a Tavant-branded Weekly Status Report presentation with accomplishments, priorities, risks, and milestones",
    "inputSchema": {
        "type": "object",
        "properties": {
            "project_name": {
                "type": "string",
                "description": "Name of the project"
            },
            "period_label": {
                "type": "string",
                "description": "Report period (e.g., 'Week Ending Dec 13, 2025')"
            },
            "accomplishments": {
                "type": "array",
                "items": {"type": "string"},
                "description": "List of key accomplishments from last period"
            },
            "priorities": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "description": {"type": "string"},
                        "owner": {"type": "string"}
                    }
                },
                "description": "List of top priorities with description and owner"
            },
            "risks": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "description": {"type": "string"},
                        "owner": {"type": "string"},
                        "target_date": {"type": "string"},
                        "status": {"type": "string"}
                    }
                },
                "description": "List of risks/issues with details"
            },
            "milestones": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "description": {"type": "string"},
                        "target_date": {"type": "string"},
                        "status": {"type": "string"}
                    }
                },
                "description": "List of key milestones"
            },
            "upcoming_milestones": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "description": {"type": "string"},
                        "target_date": {"type": "string"},
                        "owner": {"type": "string"}
                    }
                },
                "description": "List of upcoming milestones"
            },
            "contact_info": {
                "type": "string",
                "description": "Contact information for thank you slide"
            },
            "output_path": {
                "type": "string",
                "description": "Path to save the presentation"
            }
        },
        "required": ["project_name", "period_label", "accomplishments"]
    }
}
